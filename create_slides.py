"""
create_slides.py  —  Day 7: Presentation Generator
Generates reports/Presentation.pptx  (12 slides)

Prerequisites:
    pip3 install python-pptx

Run:
    python3 create_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import pandas as pd
from pathlib import Path
import xml.etree.ElementTree as ET

PROCESSED_DIR = Path("data/processed")
REPORTS_DIR   = Path("reports")
REPORTS_DIR.mkdir(exist_ok=True)

# ── brand colours ──────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x21, 0x37)
BLUE      = RGBColor(0x15, 0x65, 0xC0)
BLUE_MID  = RGBColor(0x19, 0x76, 0xD2)
LIGHT     = RGBColor(0xE3, 0xF2, 0xFD)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
AMBER     = RGBColor(0xF9, 0xA8, 0x25)
GREEN     = RGBColor(0x00, 0x89, 0x7B)
RED       = RGBColor(0xC6, 0x28, 0x28)
GREY_TEXT = RGBColor(0x54, 0x6E, 0x7A)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG  = RGBColor(0xF7, 0xF9, 0xFC)


def load_metrics():
    """Load key numbers for slide content."""
    m = {
        "top_fund": "ICICI Pru Midcap Fund",
        "top_score": "100.0",
        "best_3yr_cagr": "35.1%",
        "best_sharpe": "7.68",
        "worst_var": "-2.69%",
        "sip_growth": "169%",
        "sip_peak": "₹31,002 Cr",
        "aum_total": "₹81 Lakh Cr",
        "folios": "26.12 Cr",
        "n_funds": "40",
        "n_amcs": "10",
        "n_rows": "89,359",
        "n_tx": "32,778",
    }
    try:
        sc = pd.read_csv(PROCESSED_DIR / "fund_scorecard.csv")
        m["top_fund"]      = str(sc.iloc[0]["scheme_name"])[:45]
        m["top_score"]     = f"{sc.iloc[0]['score_100']:.0f}"
        m["best_3yr_cagr"] = f"{sc['cagr_3yr_pct'].max():.1f}%"
        m["best_sharpe"]   = f"{sc['sharpe_ratio'].max():.2f}"
    except:
        pass
    try:
        vr = pd.read_csv(PROCESSED_DIR / "var_cvar_report.csv")
        m["worst_var"] = f"{vr.iloc[0]['var_95_daily']:.3f}%"
    except:
        pass
    return m


# ── helper: set slide background colour ───────────────────────────────────

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── helper: add a filled rectangle ────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color: RGBColor, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


# ── helper: add text box ──────────────────────────────────────────────────

def add_text(slide, text, x, y, w, h,
             size=14, bold=False, italic=False,
             color: RGBColor = DARK_TEXT,
             align=PP_ALIGN.LEFT,
             wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    # clear default paragraph
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


# ── helper: add multi-line text with bullet points ────────────────────────

def add_bullets(slide, lines, x, y, w, h,
                size=12, color: RGBColor = DARK_TEXT,
                bullet_color: RGBColor = BLUE,
                line_spacing=1.2):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"▸  {line}" if not line.startswith("  ") else line
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txBox


# ── helper: add a KPI card ────────────────────────────────────────────────

def add_kpi(slide, x, y, w, h, value, label,
            bg=BLUE, fg=WHITE, label_color=None):
    add_rect(slide, x, y, w, h, bg)
    label_color = label_color or RGBColor(0xCA, 0xDC, 0xFC)
    add_text(slide, value, x+0.07, y+0.08, w-0.14, h*0.55,
             size=20, bold=True, color=fg, align=PP_ALIGN.CENTER)
    add_text(slide, label, x+0.07, y+h*0.58, w-0.14, h*0.38,
             size=9, color=label_color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════

def slide01_title(prs, m):
    """Slide 1: Title — dark navy, centered."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(sld, NAVY)

    # top accent bar
    add_rect(sld, 0, 0, 13.33, 0.12, BLUE)
    # bottom accent bar
    add_rect(sld, 0, 7.38, 13.33, 0.12, AMBER)

    # emoji + company
    add_text(sld, "📊  BLUESTOCK FINTECH", 1, 0.6, 11.33, 0.6,
             size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

    # main title
    add_text(sld, "Mutual Fund Analytics Platform",
             0.8, 1.4, 11.73, 1.2,
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sld, "End-to-End Data Engineering, ETL Pipeline & Interactive Dashboard",
             1, 2.75, 11.33, 0.7,
             size=16, italic=True, color=RGBColor(0xCA, 0xDC, 0xFC),
             align=PP_ALIGN.CENTER)

    # divider
    add_rect(sld, 3.5, 3.65, 6.33, 0.04, BLUE)

    # meta info
    meta = [
        "Data Analyst Internship  ·  Capstone Project  ·  Cohort 2025",
        "Riddhima Raj  ·  Manipal University Jaipur  ·  B.Tech CSE (IoT)",
        "Data: AMFI India · mfapi.in · NSE/BSE · June–July 2026",
    ]
    for i, line in enumerate(meta):
        add_text(sld, line, 1, 3.85 + i*0.45, 11.33, 0.4,
                 size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)

    # bottom stats strip
    stats = [("40", "Schemes"), ("10", "AMCs"), ("46K+", "NAV Rows"),
             ("32K", "Transactions"), ("₹81L Cr", "Industry AUM")]
    box_w = 2.4
    for i, (val, lbl) in enumerate(stats):
        add_kpi(sld, 0.27 + i*2.58, 6.0, box_w, 1.1,
                val, lbl, bg=BLUE_MID)


def slide02_problem(prs, m):
    """Slide 2: Problem Statement & Objectives."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, NAVY)
    add_text(sld, "Problem Statement & Objectives", 0.3, 0.1, 12, 0.45,
             size=20, bold=True, color=WHITE)

    # Left: Problems
    add_rect(sld, 0.2, 0.75, 6.2, 6.5, WHITE)
    add_rect(sld, 0.2, 0.75, 6.2, 0.4, BLUE)
    add_text(sld, "📌  5 Business Problems", 0.3, 0.8, 6.0, 0.35,
             size=13, bold=True, color=WHITE)
    problems = [
        "P1: Data fragmentation — NAV, AUM & SIP data in different formats",
        "P2: No unified risk-adjusted performance comparison tool",
        "P3: Retail investors can't track fund vs benchmark performance",
        "P4: AMCs lack demographic insight into investor behaviour patterns",
        "P5: Monthly static PDF reports take days to prepare for stakeholders",
    ]
    add_bullets(sld, problems, 0.3, 1.25, 6.0, 5.8, size=11.5, color=DARK_TEXT)

    # Right: Objectives
    add_rect(sld, 6.93, 0.75, 6.2, 6.5, WHITE)
    add_rect(sld, 6.93, 0.75, 6.2, 0.4, GREEN)
    add_text(sld, "🎯  8 Project Objectives", 7.03, 0.8, 6.0, 0.35,
             size=13, bold=True, color=WHITE)
    objectives = [
        "O1: Build automated ETL pipeline from raw AMFI data",
        "O2: Design normalised 5-table star schema in SQLite",
        "O3: Perform 17-chart EDA on NAV, AUM & SIP trends",
        "O4: Compute Sharpe, Sortino, Alpha, Beta, Max Drawdown",
        "O5: Build 4-page interactive Streamlit dashboard",
        "O6: Analyse investor behaviour across demographics",
        "O7: Compare fund returns vs Nifty 50 & Nifty 100",
        "O8: Deliver PDF report + presentation + clean GitHub repo",
    ]
    add_bullets(sld, objectives, 7.03, 1.25, 6.0, 5.8, size=11.5, color=DARK_TEXT)


def slide03_architecture(prs, m):
    """Slide 3: System Architecture — 5-layer ETL pipeline."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, NAVY)
    add_text(sld, "System Architecture — 5-Layer ETL Pipeline", 0.3, 0.1, 12, 0.45,
             size=20, bold=True, color=WHITE)

    layers = [
        (BLUE,     "LAYER 1: EXTRACT",    "mfapi.in REST API  ·  10 CSV datasets  ·  AMFI India public data  ·  NSE/BSE benchmark prices"),
        (BLUE_MID, "LAYER 2: TRANSFORM",  "Pandas cleaning  ·  Date parsing  ·  Forward-fill (NAV holidays)  ·  AMFI code validation  ·  Type coercion"),
        (GREEN,    "LAYER 3: LOAD",       "SQLite via SQLAlchemy  ·  11-table star schema  ·  89,359 rows  ·  Indexed on amfi_code + date"),
        (RGBColor(0x6A,0x1B,0x9A), "LAYER 4: ANALYSE",  "EDA notebooks  ·  CAGR / Sharpe / Sortino / Alpha / Beta  ·  VaR / CVaR  ·  Cohort  ·  HHI"),
        (AMBER,    "LAYER 5: VISUALISE",  "Streamlit 4-page dashboard  ·  20+ PNG charts  ·  KPI cards  ·  Interactive slicers  ·  Fund selector"),
    ]

    y = 0.75
    for i, (color, title, desc) in enumerate(layers):
        add_rect(sld, 0.3, y, 12.73, 0.52, color)
        add_text(sld, title, 0.45, y+0.05, 3.5, 0.4,
                 size=11, bold=True, color=WHITE)
        add_text(sld, desc,  3.9,   y+0.05, 9.0, 0.4,
                 size=10, color=RGBColor(0xE0,0xE8,0xFF) if color != AMBER else DARK_TEXT)
        if i < len(layers)-1:
            add_text(sld, "▼", 6.5, y+0.52, 0.5, 0.28,
                     size=14, color=GREY_TEXT, align=PP_ALIGN.CENTER)
        y += 0.82

    # Data flow summary
    add_rect(sld, 0.3, 6.1, 12.73, 0.42, LIGHT)
    add_text(sld, "📂  data/raw/  →  data/processed/  →  data/db/bluestock_mf.db  →  notebooks/  →  dashboard.py",
             0.5, 6.15, 12.33, 0.35, size=10.5, color=DARK_TEXT, align=PP_ALIGN.CENTER)


def slide04_database(prs, m):
    """Slide 4: Database Design — star schema overview."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, NAVY)
    add_text(sld, "Database Design — SQLite Star Schema", 0.3, 0.1, 12, 0.45,
             size=20, bold=True, color=WHITE)

    # table summary
    tables = [
        ("dim_fund",         "Dimension", "40",     "amfi_code (PK)", "40 mutual fund schemes — all metadata"),
        ("dim_date",         "Dimension", "1,826",  "date (PK)",      "Calendar 2022–2026, year/month/quarter"),
        ("fact_nav",         "Fact",      "46,000", "(amfi_code,date)","Daily NAV + computed daily return %"),
        ("fact_transactions","Fact",      "32,778", "auto",           "SIP / Lumpsum / Redemption investor events"),
        ("fact_performance", "Fact",      "40",     "amfi_code (PK)", "CAGR, Sharpe, Alpha, Beta, Max Drawdown"),
        ("fact_aum",         "Fact",      "90",     "auto",           "Quarterly AUM by AMC (₹ lakh crore)"),
        ("sip_inflows",      "Supporting","48",     "month (PK)",     "Industry monthly SIP statistics"),
        ("category_inflows", "Supporting","144",    "auto",           "Net inflows by fund category per month"),
        ("folio_count",      "Supporting","21",     "month (PK)",     "Total MF investor accounts over time"),
        ("portfolio_holdings","Supporting","322",   "auto",           "Top stock holdings per equity fund"),
        ("benchmark_indices","Supporting","8,050",  "auto",           "Nifty 50/100/Midcap150/BSESmallCap daily"),
    ]

    # Header
    cols = [1.8, 1.1, 0.8, 2.0, 4.73]
    headers = ["Table", "Type", "Rows", "Primary Key", "Description"]
    x_starts = [0.2, 2.1, 3.25, 4.1, 6.2]
    y = 0.72
    for xi, (hdr, col_w) in enumerate(zip(headers, cols)):
        add_rect(sld, x_starts[xi], y, col_w-0.05, 0.32, NAVY)
        add_text(sld, hdr, x_starts[xi]+0.05, y+0.04, col_w-0.1, 0.25,
                 size=9.5, bold=True, color=WHITE)

    for i, row in enumerate(tables):
        y = 1.06 + i * 0.53
        bg = LIGHT if i % 2 == 0 else WHITE
        for xi, (val, col_w) in enumerate(zip(row, cols)):
            add_rect(sld, x_starts[xi], y, col_w-0.05, 0.50, bg)
            clr = BLUE if xi == 0 else DARK_TEXT
            add_text(sld, val, x_starts[xi]+0.05, y+0.05, col_w-0.1, 0.42,
                     size=8.5, bold=(xi==0), color=clr)

    # Total row
    y_total = 1.06 + len(tables) * 0.53
    add_rect(sld, 0.2, y_total, 10.93, 0.32, BLUE)
    add_text(sld, f"TOTAL  ·  11 tables  ·  89,359 rows  ·  Indexed on amfi_code + date for fast BI queries",
             0.3, y_total+0.05, 10.73, 0.25,
             size=9.5, bold=True, color=WHITE)


def slide05_eda_industry(prs, m):
    """Slide 5: EDA — Industry Trends."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, BLUE)
    add_text(sld, "EDA Highlights — Indian MF Industry Trends", 0.3, 0.1, 12, 0.45,
             size=20, bold=True, color=WHITE)

    # 4 big stat boxes
    stats = [
        (NAVY,  "₹81 Lakh Cr",  "Industry AUM\n(Dec 2025)"),
        (GREEN, "₹31,002 Cr",   "Peak SIP Inflow\n(Dec 2025 — ATH)"),
        (BLUE,  "26.12 Crore",  "Total MF Folios\n(+97% vs 2022)"),
        (RGBColor(0x6A,0x1B,0x9A), "169%", "SIP Inflow Growth\n(Jan 2022 → Dec 2025)"),
    ]
    for i, (bg, val, lbl) in enumerate(stats):
        add_kpi(sld, 0.2 + i*3.25, 0.72, 3.05, 1.3, val, lbl, bg=bg)

    # Chart references
    add_rect(sld, 0.2, 2.15, 6.35, 2.5, WHITE)
    add_rect(sld, 0.2, 2.15, 6.35, 0.35, NAVY)
    add_text(sld, "📈  SIP Inflow Trend (Chart 03a)", 0.35, 2.2, 6.0, 0.28,
             size=11, bold=True, color=WHITE)
    findings = [
        "Jan 2022: ₹11,517 Cr  →  Dec 2025: ₹31,002 Cr",
        "Consistent upward trend with no single-month decline",
        "Active SIP accounts grew from 4.91 Cr to 9.35 Cr",
        "YoY growth peaked at +45.7% in 2024",
    ]
    add_bullets(sld, findings, 0.3, 2.58, 6.15, 2.0, size=11, color=DARK_TEXT)

    add_rect(sld, 6.78, 2.15, 6.35, 2.5, WHITE)
    add_rect(sld, 6.78, 2.15, 6.35, 0.35, GREEN)
    add_text(sld, "📊  AUM Growth by AMC (Chart 02)", 6.93, 2.2, 6.0, 0.28,
             size=11, bold=True, color=WHITE)
    aum_findings = [
        "SBI: ₹6.05L Cr (2022) → ₹12.5L Cr (2025) — 2× growth",
        "ICICI, HDFC follow as 2nd and 3rd largest AMCs",
        "Mirae Asset showed fastest % growth (2.4×)",
        "Top 3 AMCs control ~42% of tracked AUM",
    ]
    add_bullets(sld, aum_findings, 6.93, 2.58, 6.1, 2.0, size=11, color=DARK_TEXT)

    # Folio finding
    add_rect(sld, 0.2, 4.78, 12.93, 0.95, LIGHT)
    add_rect(sld, 0.2, 4.78, 0.06, 0.95, GREEN)
    add_text(sld, "📁  Folio Count Growth (Chart 07)  —  Total MF folios doubled from 13.26 Cr (Jan 2022) "
             "to 26.12 Cr (Dec 2025). Equity folios grew fastest, confirming shift from savings to equity investing.",
             0.35, 4.86, 12.6, 0.75, size=11, color=DARK_TEXT)

    add_rect(sld, 0.2, 5.85, 12.93, 0.95, LIGHT)
    add_rect(sld, 0.2, 5.85, 0.06, 0.95, BLUE)
    add_text(sld, "🔥  Category Heatmap (Chart 04)  —  Liquid funds dominate net inflows in absolute value "
             "(₹451,275 Cr total). Small Cap and Mid Cap show the steepest growth trajectory in FY 2024–25.",
             0.35, 5.93, 12.6, 0.75, size=11, color=DARK_TEXT)


def slide06_eda_funds(prs, m):
    """Slide 6: EDA — Fund & Investor Analysis."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, BLUE)
    add_text(sld, "EDA Highlights — Fund & Investor Analysis", 0.3, 0.1, 12, 0.45,
             size=20, bold=True, color=WHITE)

    sections = [
        (NAVY,  "🔗  Return Correlation Matrix (Chart 08)",
         ["Large Cap funds show >0.85 pairwise correlation",
          "High correlation = limited diversification benefit",
          "Small Cap funds are more independent (<0.6)",
          "Cross-category diversification is meaningful"]),
        (GREEN, "🥧  Sector Allocation — Equity Funds (Chart 09)",
         ["Banking & Financial Services is dominant sector",
          "Reflects Nifty 50 index composition bias",
          "IT, Pharma, FMCG follow as top 3 sectors",
          "Axis Bluechip most concentrated (HHI: 2064)"]),
        (RGBColor(0x6A,0x1B,0x9A), "👥  Investor Demographics (Charts 05–06)",
         ["56+ investors have highest avg SIP (₹11,575/month)",
          "Male investors: 66.5% | Female: 33.5%",
          "Punjab, Tamil Nadu, Madhya Pradesh top 3 states",
          "B30 cities: 34.1% share & growing faster than T30"]),
        (RGBColor(0x00,0x83,0x8F), "📈  NAV Trends 2022–2026 (Chart 01a/01b)",
         ["Small Cap funds outperformed all categories",
          "2023 bull run visible across all equity funds",
          "2024 mid-year correction recovered by H2 2024",
          "Indexed NAV shows 2–3× returns for top funds"]),
    ]

    positions = [(0.2, 0.75), (6.78, 0.75), (0.2, 4.0), (6.78, 4.0)]
    for (x, y), (color, title, bullets) in zip(positions, sections):
        add_rect(sld, x, y, 6.35, 3.0, WHITE)
        add_rect(sld, x, y, 6.35, 0.38, color)
        add_text(sld, title, x+0.1, y+0.06, 6.1, 0.3, size=11, bold=True, color=WHITE)
        add_bullets(sld, bullets, x+0.15, y+0.46, 6.05, 2.45, size=11, color=DARK_TEXT)


def slide07_performance(prs, m):
    """Slide 7: Performance Analytics Results."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, NAVY)
    add_text(sld, "Fund Performance Analytics — Key Results", 0.3, 0.1, 12, 0.45,
             size=20, bold=True, color=WHITE)

    # Metric formula boxes
    metrics = [
        (BLUE,  "CAGR",    f"Best 3yr: {m['best_3yr_cagr']}\nSmall Cap dominates",        "(NAV_end/NAV_start)^(1/n)-1"),
        (GREEN, "Sharpe",  f"Best: {m['best_sharpe']}\n(ICICI Pru Liquid)",                "(Rp-Rf)/σ × √252  |  Rf=6.5%"),
        (RGBColor(0x6A,0x1B,0x9A), "Alpha", "Top: SBI Small Cap\n+30.34% ann.", "OLS intercept × 252 vs Nifty 100"),
        (RED,   "Max DD",  "Worst: SBI Small Cap\n-52.57% over 1015 days",                 "min(NAV/cummax(NAV)-1)"),
    ]
    for i, (color, metric, result, formula) in enumerate(metrics):
        x = 0.2 + i * 3.27
        add_rect(sld, x, 0.72, 3.1, 2.5, WHITE)
        add_rect(sld, x, 0.72, 3.1, 0.38, color)
        add_text(sld, metric, x+0.1, 0.77, 2.9, 0.3, size=13, bold=True, color=WHITE)
        add_text(sld, result, x+0.1, 1.18, 2.9, 0.85, size=11.5, bold=True, color=DARK_TEXT)
        add_rect(sld, x+0.1, 2.15, 2.9, 0.85, LIGHT)
        add_text(sld, f"Formula: {formula}", x+0.15, 2.2, 2.8, 0.75,
                 size=8.5, italic=True, color=GREY_TEXT)

    # Scorecard table
    add_rect(sld, 0.2, 3.35, 12.93, 0.35, NAVY)
    add_text(sld, "🏆  Fund Composite Scorecard (0–100)  —  Weighted Rank: 30% CAGR + 25% Sharpe + 20% Alpha + 15% Exp Ratio + 10% Max DD",
             0.35, 3.39, 12.63, 0.28, size=9.5, bold=True, color=WHITE)

    headers2 = ["Rank", "Fund Name", "Score", "3yr CAGR", "Sharpe", "Alpha %", "Max DD %"]
    col_x    = [0.2, 0.75, 7.55, 8.55, 9.55, 10.55, 11.6]
    col_w    = [0.5, 6.75, 0.95, 0.95, 0.95, 0.95,  1.4]
    y = 3.75
    for xi, (hdr, cw) in enumerate(zip(headers2, col_w)):
        add_rect(sld, col_x[xi], y, cw-0.02, 0.3, BLUE)
        add_text(sld, hdr, col_x[xi]+0.03, y+0.04, cw-0.06, 0.22,
                 size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    scorecard_rows = [
        ("1", "ICICI Pru Midcap Fund - Regular",    "100.0", "31.77%", "1.181", "29.26%", "-21.84%"),
        ("2", "Axis Midcap Fund - Regular",           "94.2",  "35.10%", "0.999", "27.33%", "-23.10%"),
        ("3", "HDFC Mid-Cap Opportunities - Regular", "93.8",  "32.43%", "1.094", "28.27%", "-13.67%"),
        ("4", "Mirae Asset Large Cap - Regular",      "93.0",  "33.99%", "1.449", "26.80%", "-17.07%"),
        ("5", "Kotak Flexicap Fund - Regular",        "90.3",  "29.58%", "1.307", "27.33%", "-19.50%"),
    ]
    for i, row in enumerate(scorecard_rows):
        y = 4.09 + i * 0.52
        bg = LIGHT if i % 2 == 0 else WHITE
        for xi, (val, cw) in enumerate(zip(row, col_w)):
            add_rect(sld, col_x[xi], y, cw-0.02, 0.48, bg)
            clr = BLUE if xi == 0 else (GREEN if xi == 2 else DARK_TEXT)
            add_text(sld, val, col_x[xi]+0.03, y+0.06, cw-0.06, 0.36,
                     size=9, bold=(xi in [0,2]), color=clr, align=PP_ALIGN.CENTER)


def slide08_dashboard(prs, m):
    """Slide 8: Dashboard — 4 pages overview."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, RGBColor(0x00,0x89,0x7B))
    add_text(sld, "Interactive Streamlit Dashboard — 4 Pages  (Bonus B2)", 0.3, 0.1, 12, 0.45,
             size=20, bold=True, color=WHITE)

    pages = [
        (NAVY,  "🏦  Page 1\nIndustry Overview",
         "5 KPI cards · AUM trend by AMC · SIP inflow area chart · Folio growth · Monthly benchmark"),
        (BLUE,  "📈  Page 2\nFund Performance",
         "Risk/return bubble chart · Sortable scorecard table · NAV vs Nifty 50 selector · 3 slicers"),
        (GREEN, "👥  Page 3\nInvestor Analytics",
         "State bar chart · Transaction donut · Age SIP bars · Gender split · T30/B30 · Monthly volume"),
        (RGBColor(0x6A,0x1B,0x9A), "📊  Page 4\nSIP & Market Trends",
         "Dual-axis SIP + Nifty 50 · Category heatmap · Top 5 categories · YoY growth bars"),
    ]
    positions = [(0.2, 0.72), (6.78, 0.72), (0.2, 4.0), (6.78, 4.0)]
    for (x, y), (color, title, desc) in zip(positions, pages):
        add_rect(sld, x, y, 6.35, 3.0, WHITE)
        add_rect(sld, x, y, 6.35, 0.9, color)
        add_text(sld, title, x+0.15, y+0.08, 6.0, 0.8,
                 size=13, bold=True, color=WHITE)
        add_bullets(sld, desc.split(" · "), x+0.15, y+0.98, 6.05, 1.9,
                    size=10.5, color=DARK_TEXT)

    # run command
    add_rect(sld, 0.2, 7.05, 12.93, 0.35, NAVY)
    add_text(sld, "▶  Run: streamlit run dashboard.py  →  http://localhost:8501  ·  All charts interactive  ·  Data: bluestock_mf.db",
             0.4, 7.08, 12.53, 0.28, size=10, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def slide09_investor(prs, m):
    """Slide 9: Investor Analytics deep-dive."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, NAVY)
    add_text(sld, "Investor Analytics — 32,778 Transactions · 5,000 Investors · 12 States", 0.3, 0.1, 12, 0.45,
             size=18, bold=True, color=WHITE)

    # KPI strip
    inv_stats = [
        (BLUE,  "32,778",  "Total Transactions"),
        (GREEN, "5,000",   "Unique Investors"),
        (NAVY,  "12",      "States Covered"),
        (RGBColor(0x6A,0x1B,0x9A), "65.9%", "T30 City Share"),
        (RED,   "66.5%",   "Male Investors"),
    ]
    for i, (bg, val, lbl) in enumerate(inv_stats):
        add_kpi(sld, 0.2 + i*2.6, 0.72, 2.45, 1.0, val, lbl, bg=bg)

    boxes = [
        (BLUE, "📍  Geographic Insights",
         ["Punjab, Tamil Nadu, MP — top 3 states by SIP value",
          "B30 cities contribute 34.1% — growing 2.5× faster",
          "Maharashtra highest lumpsum amounts",
          "Delhi + Telangana strongest T30 SIP markets"]),
        (GREEN, "👤  Demographic Insights",
         ["56+ age group: highest avg SIP (₹11,575/month)",
          "26–35 cohort: most SIP transactions (8,063)",
          "Female investors: higher avg income proportion",
          "UPI dominant payment mode (43% of transactions)"]),
        (RGBColor(0x6A,0x1B,0x9A), "💰  Transaction Patterns",
         ["Lumpsum: ₹205.98 Cr total (avg ₹2.54L per txn)",
          "Redemption: ₹124.45 Cr (avg ₹2.51L per txn)",
          "SIP: ₹21.72 Cr (avg ₹11,018 per txn)",
          "Peak transaction months: Q4 2024, Q1 2025"]),
        (RGBColor(0x00,0x83,0x8F), "🗓️  Cohort Analysis",
         ["2024 cohort: 4,624 investors, avg SIP ₹10,997",
          "2025 cohort: 138 investors, avg SIP ₹13,505 (↑23%)",
          "Newer investors starting with higher SIP amounts",
          "ICICI Pru Bluechip most popular among 2024 cohort"]),
    ]
    positions2 = [(0.2, 1.88), (6.78, 1.88), (0.2, 4.5), (6.78, 4.5)]
    for (x, y), (color, title, bullets) in zip(positions2, boxes):
        add_rect(sld, x, y, 6.35, 2.45, WHITE)
        add_rect(sld, x, y, 6.35, 0.35, color)
        add_text(sld, title, x+0.1, y+0.06, 6.1, 0.28, size=11, bold=True, color=WHITE)
        add_bullets(sld, bullets, x+0.15, y+0.42, 6.05, 1.95, size=10.5, color=DARK_TEXT)


def slide10_advanced(prs, m):
    """Slide 10: Advanced Analytics."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, RGBColor(0x6A,0x1B,0x9A))
    add_text(sld, "Advanced Analytics — Risk Metrics, Recommender & HHI", 0.3, 0.1, 12, 0.45,
             size=20, bold=True, color=WHITE)

    sections = [
        (NAVY, "📉  Historical VaR (95%) & CVaR",
         ["Small Cap VaR: -2.69% daily (worst)",
          "Liquid Fund VaR: -0.02% daily (safest)",
          "CVaR always worse than VaR — captures tail risk",
          "All 40 funds computed — var_cvar_report.csv"]),
        (BLUE, "📈  Rolling 90-Day Sharpe",
         ["Shows Sharpe evolution over time for 5 funds",
          "2024 market correction visible as Sharpe dip",
          "Recovery in H2 2024 confirms data validity",
          "chart: 13_rolling_sharpe.png"]),
        (GREEN, "🤖  Fund Recommender (recommender.py)",
         ["Low risk  → ICICI Pru Liquid (Sharpe 7.68)",
          "Moderate → Mirae Asset Large Cap (Sharpe 1.06)",
          "High risk → ICICI Pru Midcap (Sharpe 0.95)",
          "CLI: python3 recommender.py --risk High"]),
        (RGBColor(0x00,0x83,0x8F), "🏭  Sector HHI Concentration",
         ["HHI = Σ(weight_i²) per fund portfolio",
          "Axis Bluechip most concentrated (HHI: 2064, IT)",
          "ABSL Small Cap 2nd (HHI: 2007, Pharma)",
          "Mid Cap funds show lower HHI = better diversity"]),
    ]

    positions = [(0.2, 0.72), (6.78, 0.72), (0.2, 3.9), (6.78, 3.9)]
    for (x, y), (color, title, bullets) in zip(positions, sections):
        add_rect(sld, x, y, 6.35, 2.95, WHITE)
        add_rect(sld, x, y, 6.35, 0.38, color)
        add_text(sld, title, x+0.1, y+0.06, 6.1, 0.3, size=11, bold=True, color=WHITE)
        add_bullets(sld, bullets, x+0.15, y+0.45, 6.05, 2.4, size=11, color=DARK_TEXT)

    add_rect(sld, 0.2, 7.0, 12.93, 0.38, LIGHT)
    add_rect(sld, 0.2, 7.0, 0.06, 0.38, RGBColor(0x6A,0x1B,0x9A))
    add_text(sld, "Deliverables: Advanced_Analytics.ipynb  ·  var_cvar_report.csv  ·  cohort_analysis.csv  ·  sip_continuity.csv  ·  sector_hhi.csv  ·  recommender.py",
             0.35, 7.05, 12.58, 0.28, size=9, color=GREY_TEXT)


def slide11_findings(prs, m):
    """Slide 11: 10 Key Findings."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, LIGHT_BG)
    add_rect(sld, 0, 0, 13.33, 0.6, AMBER)
    add_text(sld, "10 Key Business Findings", 0.3, 0.1, 12, 0.45,
             size=22, bold=True, color=DARK_TEXT)

    findings = [
        ("01", "SIP inflows grew 169% in 4 years — ₹11,517 Cr (2022) to ₹31,002 Cr (Dec 2025 ATH)", BLUE),
        ("02", "SBI MF leads AUM at ₹12.5L Cr — 2× from 2022; top 3 AMCs control 42% of tracked AUM", NAVY),
        ("03", "Small Cap funds: highest 3yr CAGR (20–23%) but worst VaR (-2.69%) and drawdown (-52%)", RED),
        ("04", "ICICI Pru Midcap scored 100/100 on composite scorecard across 5 weighted metrics", GREEN),
        ("05", "MF folios doubled — 13.26 Cr → 26.12 Cr — confirming mass retail equity adoption", BLUE),
        ("06", "B30 cities: 34.1% investment share and growing 2.5× faster than T30 metros", GREEN),
        ("07", "56+ investors have highest avg SIP (₹11,575/month) — retirement planning effect", NAVY),
        ("08", "Large Cap fund returns >0.85 correlated — within-category diversification is limited", BLUE),
        ("09", "Axis Bluechip most portfolio-concentrated (HHI 2064) — heavy IT sector exposure", RED),
        ("10", "Direct plans: avg 0.7–0.9% lower expense ratio — compounds to significant advantage", GREEN),
    ]

    col1 = findings[:5]
    col2 = findings[5:]
    for col_idx, col_findings in enumerate([col1, col2]):
        x = 0.2 + col_idx * 6.65
        for i, (num, text, color) in enumerate(col_findings):
            y = 0.72 + i * 1.3
            add_rect(sld, x, y, 6.35, 1.15, WHITE)
            add_rect(sld, x, y, 0.45, 1.15, color)
            add_text(sld, num, x+0.02, y+0.3, 0.41, 0.55,
                     size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(sld, text, x+0.55, y+0.15, 5.7, 0.9,
                     size=10.5, color=DARK_TEXT)


def slide12_thankyou(prs, m):
    """Slide 12: Thank You — dark navy close."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sld, NAVY)
    add_rect(sld, 0, 0, 13.33, 0.12, BLUE)
    add_rect(sld, 0, 7.38, 13.33, 0.12, AMBER)

    add_text(sld, "Thank You", 1, 1.0, 11.33, 1.1,
             size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sld, "Bluestock MF Analytics Platform", 1, 2.3, 11.33, 0.6,
             size=18, color=AMBER, align=PP_ALIGN.CENTER)
    add_rect(sld, 3.5, 3.05, 6.33, 0.05, BLUE)

    add_text(sld, "Riddhima Raj  ·  Data Analyst Intern  ·  Cohort 2025", 1, 3.25, 11.33, 0.45,
             size=13, color=RGBColor(0xCA,0xDC,0xFC), align=PP_ALIGN.CENTER)
    add_text(sld, "B.Tech CSE (IoT & Intelligent Systems)  ·  Manipal University Jaipur  ·  CGPA 9.10",
             1, 3.75, 11.33, 0.4, size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
    add_text(sld, "🔗  github.com/riddhima16/bluestock-mf-capstone",
             1, 4.3, 11.33, 0.4, size=12, color=AMBER, align=PP_ALIGN.CENTER)

    # tech stack badges
    techs = ["Python 3.11", "Pandas", "SQLite", "Plotly", "Streamlit", "Scipy", "Git"]
    total_w = len(techs) * 1.6 + (len(techs)-1) * 0.15
    start_x = (13.33 - total_w) / 2
    for i, tech in enumerate(techs):
        add_rect(sld, start_x + i*1.75, 5.2, 1.55, 0.4, BLUE_MID)
        add_text(sld, tech, start_x + i*1.75 + 0.05, 5.24, 1.45, 0.32,
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sld, "⚠  Disclaimer: All AMFI codes and real AUM/SIP figures are from public AMFI India data. "
             "NAV values are anchored to real mfapi.in values. Investor transactions are synthetic. "
             "This project is for educational purposes only and does not constitute financial advice.",
             0.5, 5.85, 12.33, 0.9, size=8.5, italic=True,
             color=GREY_TEXT, align=PP_ALIGN.CENTER)

    # Bluestock branding
    add_text(sld, "📊 BLUESTOCK FINTECH  ·  MF Analytics Platform  ·  Data Analyst Internship Capstone 2025",
             0.5, 6.9, 12.33, 0.35, size=9, color=RGBColor(0x54,0x6E,0x7A),
             align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("  create_slides.py — Day 7: Presentation Generator")
    print("=" * 60)
    print("\nPre-requisite: pip3 install python-pptx")

    m = load_metrics()
    print(f"\n  Key metrics loaded: top fund = {m['top_fund'][:40]}")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("\n  Building 12 slides...")
    slide01_title(prs, m);          print("    ✓ Slide  1: Title")
    slide02_problem(prs, m);        print("    ✓ Slide  2: Problem & Objectives")
    slide03_architecture(prs, m);   print("    ✓ Slide  3: System Architecture")
    slide04_database(prs, m);       print("    ✓ Slide  4: Database Design")
    slide05_eda_industry(prs, m);   print("    ✓ Slide  5: EDA — Industry Trends")
    slide06_eda_funds(prs, m);      print("    ✓ Slide  6: EDA — Fund & Investor")
    slide07_performance(prs, m);    print("    ✓ Slide  7: Performance Analytics")
    slide08_dashboard(prs, m);      print("    ✓ Slide  8: Dashboard Overview")
    slide09_investor(prs, m);       print("    ✓ Slide  9: Investor Analytics")
    slide10_advanced(prs, m);       print("    ✓ Slide 10: Advanced Analytics")
    slide11_findings(prs, m);       print("    ✓ Slide 11: Key Findings")
    slide12_thankyou(prs, m);       print("    ✓ Slide 12: Thank You")

    out_path = REPORTS_DIR / "Presentation.pptx"
    prs.save(str(out_path))
    print(f"\n  ✓ Saved: {out_path}")
    print(f"  Open in PowerPoint or Google Slides to view.")
    print(f"\n  Day 7 complete — all deliverables ready for submission!")


if __name__ == "__main__":
    main()